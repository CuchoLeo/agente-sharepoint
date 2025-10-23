# -*- coding: utf-8 -*-
# Importaciones necesarias para el funcionamiento de la aplicación
import os  # Para interactuar con el sistema operativo (leer variables de entorno)
import io  # Para manejar flujos de bytes en memoria (archivos)
import json  # Para trabajar con archivos JSON (el log de indexación)
import time  # Para añadir pausas y registrar marcas de tiempo
import numpy as np  # Para operaciones numéricas, especialmente con los vectores de embeddings
import requests  # Para realizar peticiones HTTP a las APIs de Microsoft Graph
from fastapi import FastAPI, Form, UploadFile, File  # El framework web para construir la API
from fastapi.responses import HTMLResponse, JSONResponse  # Tipos de respuesta para la API
from fastapi.staticfiles import StaticFiles  # <-- MODIFICACIÓN 1: Importar StaticFiles
from dotenv import load_dotenv  # Para cargar variables de entorno desde un archivo .env
from openai import OpenAI  # Cliente oficial de OpenAI para interactuar con GPT y modelos de embedding
from msal import ConfidentialClientApplication  # Librería de Microsoft para la autenticación OAuth2 (MSAL)
from docx import Document  # Para leer archivos .docx de Word
from PyPDF2 import PdfReader  # Para leer archivos .pdf
from pptx import Presentation  # Para leer archivos .pptx de PowerPoint
from openpyxl import load_workbook  # Para leer archivos .xlsx de Excel
from pdf2image import convert_from_bytes  # Para convertir páginas de PDF a imágenes (necesario para OCR)
from pytesseract import image_to_string  # Para realizar OCR (Reconocimiento Óptico de Caracteres) en imágenes

# ======================================================
# CONFIGURACIÓN INICIAL
# ======================================================
load_dotenv()  # Carga las variables del archivo .env en el entorno del sistema
app = FastAPI()  # Crea una instancia de la aplicación FastAPI

# <-- MODIFICACIÓN 2: Montar el directorio estático
# Monta el directorio actual ('.') en la ruta '/static' para servir archivos
app.mount("/static", StaticFiles(directory="."), name="static")

# --- Variables de Entorno y Configuración ---
# Lee las credenciales y configuraciones desde las variables de entorno
OPENAI_KEY = os.getenv("OPENAI_API_KEY")
M365_TENANT_ID = os.getenv("M365_TENANT_ID")
M365_CLIENT_ID = os.getenv("M365_CLIENT_ID")
M365_CLIENT_SECRET = os.getenv("M365_CLIENT_SECRET")
DRIVE_ID = "b!rwJYlL5mUEuss_GKf1rBSELmRer4qdhIv5HsDY3tNOgeIMFwY0seSr-QtklnZZRh"  # ID del Drive de SharePoint a utilizar

# --- Clientes y Constantes ---
# Inicializa los clientes de OpenAI para la generación de texto y embeddings
client = OpenAI(api_key=OPENAI_KEY)
emb_client = OpenAI(api_key=OPENAI_KEY)
EMBED_MODEL = "text-embedding-3-small"  # Modelo específico de OpenAI para crear embeddings
INDEX_FILE = "index.npy"  # Nombre del archivo que almacenará el índice vectorial
INDEX_LOG_FILE = "index_log.json"  # Nombre del archivo que registrará los metadatos de los archivos indexados
MAX_FILES_PER_RUN = 10  # Límite de archivos a procesar en cada ejecución de indexación para no sobrecargar la memoria

# ======================================================
# DICCIONARIO DE CARPETAS INDEXABLES
# ======================================================
# --- *** MODIFICACIÓN *** ---
# Se ha limitado el diccionario para que solo incluya la carpeta "PRUEBAS"
SHAREPOINT_FOLDERS = {
    "PRUEBAS": "01OUATYK74IICILGLVQBDIXLFNJN6ESTI6"
}
# --- *** FIN DE LA MODIFICACIÓN *** ---

# --- Variables Globales en Memoria ---
# Se usarán para mantener el estado del índice mientras la aplicación se ejecuta
doc_chunks = []  # Lista que contendrá los fragmentos de texto y sus embeddings
index_log = {}  # Diccionario que contendrá los metadatos de los archivos de SharePoint indexados

# ======================================================
# LÓGICA DE EXTRACCIÓN DE TEXTO (REUTILIZABLE)
# ======================================================
def extract_text_from_bytes(file_bytes, file_name):
    """
    Función centralizada que recibe los bytes de un archivo y extrae su contenido de texto.
    Soporta PDF (con fallback a OCR), DOCX, XLSX y PPTX.
    """
    text = ""  # Inicializa la variable que almacenará el texto extraído
    try:
        # Lógica para archivos PDF
        if file_name.lower().endswith(".pdf"):
            try:
                # Intenta leer el texto directamente del PDF
                reader = PdfReader(io.BytesIO(file_bytes))
                text = "\n".join(page.extract_text() or "" for page in reader.pages)
                if not text.strip(): raise ValueError("PDF vacío, usar OCR")
            except Exception:
                # Si falla o el PDF es una imagen, convierte cada página a imagen y aplica OCR
                text = "\n".join(image_to_string(p, lang='spa') for p in convert_from_bytes(file_bytes))
        
        # Lógica para archivos DOCX (Word)
        elif file_name.lower().endswith(".docx"):
            doc = Document(io.BytesIO(file_bytes))
            text = "\n".join(p.text for p in doc.paragraphs)
        
        # Lógica para archivos XLSX (Excel)
        elif file_name.lower().endswith((".xlsx", ".xls")):
            wb = load_workbook(io.BytesIO(file_bytes), data_only=True) # data_only=True para obtener el valor de las celdas, no las fórmulas
            for sheet in wb.worksheets:
                sheet_text = "\n".join(" | ".join(str(c.value) if c.value is not None else "" for c in row) for row in sheet.iter_rows())
                if sheet_text.strip(): text += f"\n--- Hoja: {sheet.title} ---\n{sheet_text}"
        
        # Lógica para archivos PPTX (PowerPoint)
        elif file_name.lower().endswith(".pptx"):
            prs = Presentation(io.BytesIO(file_bytes))
            for i, slide in enumerate(prs.slides, 1):
                slide_text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                if slide_text.strip(): text += f"\n--- Diapositiva {i} ---\n{slide_text}"
        
        # Para otros tipos de archivo (ej. .txt), intenta decodificar como texto plano
        else:
            text = file_bytes.decode("utf-8", errors="ignore")
    
    # Captura cualquier excepción durante la lectura para evitar que la aplicación se caiga
    except Exception as e:
        text = f"[Error leyendo {file_name}: {e}]"
    
    # Limita la cantidad de texto para controlar el uso de memoria y costos de API
    return text[:20000]

# ======================================================
# LÓGICA DE AUTENTICACIÓN Y ACCESO A SHAREPOINT
# ======================================================
def get_graph_token():
    """Obtiene un token de acceso de Microsoft Graph usando las credenciales de la aplicación."""
    authority = f"https://login.microsoftonline.com/{M365_TENANT_ID}"
    app_msal = ConfidentialClientApplication(M365_CLIENT_ID, authority=authority, client_credential=M365_CLIENT_SECRET)
    token_result = app_msal.acquire_token_for_client(scopes=["https://graph.microsoft.com/.default"])
    if "access_token" not in token_result:
        print(f"❌ Error al obtener token: {token_result.get('error_description', 'No description')}")
        return None
    return token_result["access_token"]

def list_sharepoint_files(folder_id=None):
    """Obtiene una lista recursiva de todos los archivos dentro de una carpeta de SharePoint (o desde la raíz)."""
    token = get_graph_token()
    if not token:
        return []  # Si no se obtiene el token, retorna una lista vacía
        
    headers = {"Authorization": f"Bearer {token}"}
    url = f"https://graph.microsoft.com/v1.0/drives/{DRIVE_ID}/items/{folder_id}/children" if folder_id else f"https://graph.microsoft.com/v1.0/drives/{DRIVE_ID}/root/children"
    
    try:
        # Realiza la petición a la API de Microsoft Graph
        r = requests.get(url, headers=headers)
        r.raise_for_status()  # Lanza una excepción si la respuesta es un error HTTP (ej. 401, 403)
        items = r.json().get("value", [])
        files = []
        for item in items:
            if "file" in item:  # Si es un archivo, lo añade a la lista
                files.append(item)
            elif "folder" in item:  # Si es una carpeta, llama a esta misma función para explorar su contenido
                files.extend(list_sharepoint_files(item["id"]))
        return files
    except requests.exceptions.RequestException as e:
        # Captura errores de red o de autenticación y los muestra en la consola
        print(f"❌ Error al conectar con SharePoint: {e}")
        return []

def read_sharepoint_file_content(file_id):
    """Descarga y devuelve el contenido en bytes de un archivo específico de SharePoint."""
    token = get_graph_token()
    if not token: return None
    headers = {"Authorization": f"Bearer {token}"}
    try:
        # Primero obtiene los metadatos del archivo, que incluyen la URL de descarga
        meta_res = requests.get(f"https://graph.microsoft.com/v1.0/drives/{DRIVE_ID}/items/{file_id}", headers=headers)
        meta_res.raise_for_status()
        download_url = meta_res.json().get("@microsoft.graph.downloadUrl")
        if not download_url: return None
        # Descarga el contenido del archivo desde la URL obtenida
        file_res = requests.get(download_url)
        file_res.raise_for_status()
        return file_res.content
    except requests.exceptions.RequestException: return None

# ======================================================
# LÓGICA DE CARGA Y GUARDADO DEL ÍNDICE
# ======================================================

def save_index():
    """Guarda el índice vectorial y el log de metadatos en archivos locales."""
    try:
        # Guarda los embeddings y fragmentos de texto
        np.save(INDEX_FILE, doc_chunks, allow_pickle=True)
        
        # Guarda los metadatos en un archivo JSON
        with open(INDEX_LOG_FILE, 'w', encoding='utf-8') as f:
            json.dump(index_log, f, ensure_ascii=False, indent=4)
        
        print(f"✅ Índice y log guardados con éxito en '{INDEX_FILE}' y '{INDEX_LOG_FILE}'.")
        
    except Exception as e:
        print(f"❌ ¡ERROR CRÍTICO AL GUARDAR EL ÍNDICE!: {e}")
        print("Verifica los permisos de escritura en el directorio de la aplicación.")

def load_index():
    """Carga el índice y el log desde los archivos locales a la memoria al iniciar la aplicación."""
    global doc_chunks, index_log
    try:
        if os.path.exists(INDEX_FILE): 
            doc_chunks = np.load(INDEX_FILE, allow_pickle=True).tolist()
        if os.path.exists(INDEX_LOG_FILE):
            with open(INDEX_LOG_FILE, 'r', encoding='utf-8') as f: 
                index_log = json.load(f)
        print(f"📂 Índice cargado. {len(doc_chunks)} fragmentos y {len(index_log)} archivos en log.")
    except Exception as e:
        print(f"⚠️ Error al cargar el índice (puede ser la primera ejecución): {e}")
        doc_chunks = []
        index_log = {}

# Carga el índice en memoria en el momento que Uvicorn importa la app.
load_index() 


# ======================================================
# MOTOR DE INDEXACIÓN
# ======================================================
def build_sharepoint_index(folder_id=None, full_reindex=False):
    """Construye o actualiza el índice vectorial a partir de los archivos de SharePoint."""
    global doc_chunks, index_log
    # Si es una reindexación completa, borra los datos en memoria
    if full_reindex: 
        print("🌀 Solicitada reindexación completa. Limpiando índice y log...")
        doc_chunks, index_log = [], {}
    
    # Obtiene la lista actual de todos los archivos en SharePoint
    all_files_from_sp = list_sharepoint_files(folder_id)
    # Crea un mapa de los archivos actuales para una búsqueda eficiente
    current_files_map = { (f.get("parentReference", {}).get("path", "").split("root:")[-1] or "/") + f["name"]: {"id": f["id"], "last_modified": f.get("lastModifiedDateTime"), "name": f["name"]} for f in all_files_from_sp if "file" in f }
    
    # Determina qué archivos necesitan ser procesados (son nuevos o han sido modificados)
    files_to_process = [data for path, data in current_files_map.items() if path not in index_log or index_log[path]["last_modified"] != data["last_modified"]]
    
    # Determina qué archivos han sido eliminados de SharePoint para quitarlos del índice
    deleted_paths = set(index_log.keys()) - set(current_files_map.keys())
    if deleted_paths:
        print(f"🗑️ Se detectaron {len(deleted_paths)} archivos eliminados. Limpiando del índice...")
        ids_to_remove = {index_log[path]["id"] for path in deleted_paths}
        doc_chunks = [c for c in doc_chunks if c.get("source_id") not in ids_to_remove]
        for path in deleted_paths:
            del index_log[path] # Elimina también del log
            
    print(f"🔎 Se encontraron {len(files_to_process)} archivos nuevos o modificados para procesar.")
    
    # Itera sobre los archivos a procesar (limitado por MAX_FILES_PER_RUN)
    processed_count = 0
    for file_data in files_to_process[:MAX_FILES_PER_RUN]:
        print(f"📄 Procesando: {file_data['name']}...")
        file_bytes = read_sharepoint_file_content(file_data["id"])
        if not file_bytes: 
            print(f"⚠️ No se pudo descargar {file_data['name']}. Saltando.")
            continue
        
        text = extract_text_from_bytes(file_bytes, file_data["name"])
        if not text.strip() or text.startswith("[Error"):
            print(f"⚠️ No se pudo extraer texto de {file_data['name']}. Saltando.")
            continue
        
        # Elimina los fragmentos antiguos del archivo que se va a reindexar
        doc_chunks = [c for c in doc_chunks if c.get("source_id") != file_data["id"]]
        
        # Divide el texto en fragmentos (chunks) y genera sus embeddings
        for chunk_text in [text[j:j+1000] for j in range(0, len(text), 1000)]:
            try:
                # Llama a la API de OpenAI para crear el embedding
                vec = emb_client.embeddings.create(model=EMBED_MODEL, input=chunk_text).data[0].embedding
                # Añade el fragmento, su embedding y metadatos a la lista en memoria
                doc_chunks.append({"text": chunk_text, "embedding": np.array(vec), "source": file_data["name"], "source_id": file_data["id"]})
            except Exception as e: print(f"❌ Error en embedding para {file_data['name']}: {e}")
        
        # Actualiza el log para este archivo específico
        path_key = (file_data.get("parentReference", {}).get("path", "").split("root:")[-1] or "/") + file_data["name"]
        index_log[path_key] = file_data
        processed_count += 1
        time.sleep(0.2) # Pequeña pausa para no saturar la API
        
    # Guarda los cambios en disco
    save_index() 
    return f"Se detectaron {len(all_files_from_sp)} archivos. Se procesaron {processed_count} en este lote."

# ======================================================
# MOTOR IA Y BÚSQUEDA (LÓGICA HÍBRIDA) - CORREGIDO
# ======================================================
def search_similar_chunks(query, top_k=3):
    """Busca los fragmentos de texto más relevantes para una consulta usando similitud de coseno."""
    if not doc_chunks: return []
    # Crea el embedding de la pregunta del usuario
    q_emb = np.array(emb_client.embeddings.create(model=EMBED_MODEL, input=query).data[0].embedding)
    # Calcula la similitud (producto punto) entre la pregunta y todos los chunks del índice
    sims = sorted([(np.dot(q_emb, d["embedding"]), d) for d in doc_chunks], key=lambda x: x[0], reverse=True)
    # Devuelve el texto de los fragmentos más similares
    return [d["text"] for _, d in sims[:top_k]]

def ask_openai(question):
    """Función principal que procesa la pregunta del usuario y genera una respuesta."""
    
    # --- CORRECCIÓN: Detección de intención para listar archivos ---
    list_keywords = [
        'archivos indexad',  # Detectará 'indexados', 'indexados?', 'inndexados', etc.
        'documentos conoces',
        'lista de archivos',
        'lista de documentos',
        'que archivos tienes', # Sin tilde
        'que archivos fueron'  # Sin tilde
    ]
    
    # Comprueba si alguna de estas palabras clave flexibles está en la pregunta
    if any(keyword in question.lower() for keyword in list_keywords):
        print("▶️ Intención detectada: Listar archivos.")
        # Obtiene una lista única y ordenada de los nombres de los archivos
        unique_sources = sorted(list(set(chunk.get('source', 'Desconocido') for chunk in doc_chunks)))
        if not unique_sources:
            return "Actualmente no tengo ningún archivo indexado en mi memoria."
        # Formatea la lista de archivos para mostrarla
        file_list_str = "\n- ".join(unique_sources)
        return f"Tengo conocimiento sobre los siguientes documentos:\n\n- {file_list_str}"
    # --- FIN DE LA CORRECCIÓN ---

    # --- Lógica de búsqueda híbrida ---
    context_text, source_found = "", None
    unique_sources = sorted(list(set(chunk['source'] for chunk in doc_chunks)))
    # 1. Búsqueda por nombre de archivo: si la pregunta menciona un documento, usa ese como contexto
    for source_name in unique_sources:
        if os.path.splitext(source_name)[0].lower().replace("_", " ") in question.lower():
            source_found = source_name
            context_text = "\n\n".join([chunk['text'] for chunk in doc_chunks if chunk['source'] == source_name]).strip()
            break
    # 2. Búsqueda semántica: si no se mencionó un archivo, busca los fragmentos más similares
    if not source_found:
        context_text = "\n\n".join(search_similar_chunks(question)).strip()

    if not context_text: context_text = "No se encontró información relevante."
    
    # Construye el prompt para la IA, indicando que debe basarse en el contexto
    system_prompt = "Eres OVAL Agente IA, un asistente documental. Responde de forma concisa y profesional, basándote en el contexto proporcionado."
    if source_found: system_prompt += f" La información fue extraída del documento: {source_found}."

    # Define los mensajes que se enviarán a la API de OpenAI
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "system", "content": f"Contexto extraído de documentos:\n{context_text}"},
        {"role": "user", "content": question}
    ]

    # Llama a la API de OpenAI para generar la respuesta
    response = client.chat.completions.create(model="gpt-4o-mini", messages=messages, temperature=0.2, max_tokens=1000)
    return response.choices[0].message.content

# ======================================================
# ENDPOINTS DE LA API (CON CARGA DE ARCHIVOS)
# ======================================================
@app.get("/", response_class=HTMLResponse)
def index():
    """Sirve la página web principal (la interfaz del chat)."""
    # <-- MODIFICACIÓN 3: HTML del encabezado actualizado para incluir la imagen
    return """
<!DOCTYPE html>
<html lang="es">
<head><meta charset="utf-8"><title>OVAL Agente IA</title><script src="https://cdn.tailwindcss.com"></script></head>
<body class="bg-gray-100 flex flex-col items-center min-h-screen p-4 font-sans">
    <div class="bg-white shadow-lg rounded-2xl w-full max-w-3xl flex flex-col h-[95vh]">
        
        <div class="p-4 border-b bg-gray-600 text-white rounded-t-2xl flex justify-between items-center">
            
            <div class="flex items-center gap-3">
                <img src="/static/oval.jpg" alt="Logo OVAL" class="h-10 w-auto rounded-lg">
                
                <div>
                    <h2 class="text-xl font-bold">OVAL Agente IA</h2>
                    <p class="text-xs opacity-80">Integrado con Microsoft 365 | Demo OVAL</p>
                </div>
            </div>

            <button id="settingsBtn" class="p-2 rounded-full hover:bg-orange-700 transition"><svg xmlns="http://www.w3.org/2000/svg" class="h-6 w-6" fill="none" viewBox="0 0 24 24" stroke="currentColor"><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M10.325 4.317c.426-1.756 2.924-1.756 3.35 0a1.724 1.724 0 002.573 1.066c1.543-.94 3.31.826 2.37 2.37a1.724 1.724 0 001.065 2.572c1.756.426 1.756 2.924 0 3.35a1.724 1.724 0 00-1.066 2.573c.94 1.543-.826 3.31-2.37 2.37a1.724 1.724 0 00-2.572 1.065c-.426 1.756-2.924 1.756-3.35 0a1.724 1.724 0 00-2.573-1.066c-1.543.94-3.31-.826-2.37-2.37a1.724 1.724 0 00-1.065-2.572c-1.756-.426-1.756-2.924 0-3.35a1.724 1.724 0 001.066-2.573c-.94-1.543.826-3.31 2.37-2.37.996.608 2.296.07 2.572-1.065z" /><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M15 12a3 3 0 11-6 0 3 3 0 016 0z" /></svg></button>
        </div>
        <div id="settingsPanel" class="p-4 border-b bg-gray-50 hidden">
            <h3 class="font-bold mb-2">Mantenedor de Índice (SharePoint)</h3>
            <div class="flex items-center gap-4"><select id="folderSelect" class="flex-1 border rounded-lg p-2 text-sm focus:ring-orange-500 focus:border-orange-500"><option value="">-- Indexar todo (Raíz) --</option></select><button id="reindexBtn" class="bg-blue-600 hover:bg-blue-700 text-white font-semibold px-4 py-2 rounded-lg transition text-sm">Indexar</button><button id="reindexFullBtn" class="bg-red-600 hover:bg-red-700 text-white font-semibold px-4 py-2 rounded-lg transition text-sm">Forzar Todo</button></div>
            <div id="reindexStatus" class="text-xs text-gray-500 mt-2"></div>
        </div>
        <div id="chatBox" class="flex-1 overflow-y-auto p-4 space-y-4 bg-gray-100"></div>
        <form id="uploadForm" class="p-4 border-t flex items-center gap-2 bg-gray-50">
            <input id="fileInput" name="file" type="file" accept=".pdf,.docx,.xlsx,.pptx" class="flex-1 text-sm text-slate-500 file:mr-4 file:py-2 file:px-4 file:rounded-full file:border-0 file:text-sm file:font-semibold file:bg-orange-50 file:text-orange-700 hover:file:bg-orange-100">
            <button type="submit" class="bg-gray-700 hover:bg-gray-800 text-white font-semibold px-4 py-2 rounded-lg transition text-sm">Cargar e Indexar</button>
        </form>
        <form id="chatForm" class="p-4 border-t flex items-center gap-2 bg-white"><input id="question" name="question" class="flex-1 border rounded-lg p-2 text-sm focus:ring-orange-500 focus:border-orange-500" placeholder="Escribe tu pregunta aquí..." required><button type="submit" class="bg-orange-600 hover:bg-orange-700 text-white font-semibold px-4 py-2 rounded-lg transition">Enviar</button></form>
    </div>
<script>
// Lógica del frontend (JavaScript) que se ejecuta en el navegador del usuario
document.addEventListener('DOMContentLoaded', () => {
    // Referencias a los elementos del DOM
    const chatForm = document.getElementById('chatForm'), uploadForm = document.getElementById('uploadForm'), chatBox = document.getElementById('chatBox'), settingsBtn = document.getElementById('settingsBtn'), settingsPanel = document.getElementById('settingsPanel'), folderSelect = document.getElementById('folderSelect'), reindexBtn = document.getElementById('reindexBtn'), reindexFullBtn = document.getElementById('reindexFullBtn'), reindexStatus = document.getElementById('reindexStatus');
    
    // Al cargar la página, pide al backend la lista de carpetas de SharePoint y las añade al menú desplegable
    fetch('/get_folders').then(res => res.json()).then(folders => { Object.keys(folders).sort().forEach(name => { const opt = document.createElement('option'); opt.value = name; opt.textContent = name; folderSelect.appendChild(opt); }); });
    
    // Muestra u oculta el panel de configuración al hacer clic en el engranaje
    settingsBtn.addEventListener('click', () => settingsPanel.classList.toggle('hidden'));
    
    // Función para manejar las peticiones de reindexación
    const handleReindex = (isFull) => { const folder = folderSelect.value, url = new URL(window.location.origin + '/reindex'); if (folder) url.searchParams.append('folder', folder); if (isFull) url.searchParams.append('full', 'true'); reindexStatus.textContent = 'Indexando, espera...'; reindexBtn.disabled = reindexFullBtn.disabled = true; fetch(url).then(res => res.json()).then(data => { reindexStatus.textContent = `✅ ${data.message}`; }).catch(() => { reindexStatus.textContent = '❌ Error.'; }).finally(() => { reindexBtn.disabled = reindexFullBtn.disabled = false; }); };
    
    // Asigna la función a los botones de indexación
    reindexBtn.addEventListener('click', () => handleReindex(false));
    reindexFullBtn.addEventListener('click', () => handleReindex(true));
    
    // Lógica para el formulario de carga de archivos
    uploadForm.addEventListener('submit', async (e) => { e.preventDefault(); const fileInput = document.getElementById('fileInput'); if (!fileInput.files.length) return; const formData = new FormData(); formData.append('file', fileInput.files[0]); const statusBubble = document.createElement('div'); statusBubble.className = 'text-center text-sm text-gray-500 p-2'; statusBubble.innerText = `Indexando ${fileInput.files[0].name}...`; chatBox.appendChild(statusBubble); chatBox.scrollTop = chatBox.scrollHeight; try { const res = await fetch('/upload', { method: 'POST', body: formData }); const result = await res.json(); statusBubble.className = res.ok ? 'bg-green-100 text-green-800 p-3 rounded-lg text-sm text-center' : 'bg-red-100 text-red-800 p-3 rounded-lg text-sm text-center'; statusBubble.textContent = result.message; } catch (error) { statusBubble.className = 'bg-red-100 text-red-800 p-3 rounded-lg text-sm text-center'; statusBubble.textContent = '❌ Error de conexión.'; } uploadForm.reset(); });
    
    // Lógica para el formulario de chat
    chatForm.addEventListener('submit', async (e) => { e.preventDefault(); const formData = new FormData(chatForm), question = formData.get('question'); if (!question.trim()) return; const userBubble = document.createElement('div'); userBubble.className = 'bg-orange-100 p-3 rounded-lg self-end text-right text-gray-800 max-w-prose ml-auto'; userBubble.innerText = question; chatBox.appendChild(userBubble); chatForm.reset(); const thinking = document.createElement('div'); thinking.className = 'text-sm text-gray-400 italic self-start'; thinking.innerText = 'OVAL Agente IA está pensando...'; chatBox.appendChild(thinking); chatBox.scrollTop = chatBox.scrollHeight; const res = await fetch('/ask', { method: 'POST', body: formData }); const html = await res.text(); thinking.remove(); const botBubble = document.createElement('div'); botBubble.className = 'bg-gray-200 p-3 rounded-lg text-gray-800 max-w-prose mr-auto'; botBubble.innerHTML = html; chatBox.appendChild(botBubble); chatBox.scrollTop = chatBox.scrollHeight; });
});
</script>
</body></html>
"""

@app.get("/get_folders", response_class=JSONResponse)
def get_folders():
    """Devuelve el diccionario de carpetas de SharePoint para el frontend."""
    return SHAREPOINT_FOLDERS

@app.post("/ask", response_class=HTMLResponse)
async def ask(question: str = Form(...)):
    """Recibe la pregunta del usuario, la procesa y devuelve la respuesta de la IA."""
    return f"<p class='whitespace-pre-line'>{ask_openai(question)}</p>"

@app.post("/upload", response_class=JSONResponse)
async def upload_file(file: UploadFile = File(...)):
    """Maneja la subida de un archivo local, lo procesa y lo añade al índice."""
    global doc_chunks, index_log
    file_bytes = await file.read() # Lee el contenido del archivo subido
    text = extract_text_from_bytes(file_bytes, file.filename) # Extrae el texto
    
    # Si no se pudo extraer texto, devuelve un error
    if not text.strip() or text.startswith("[Error"): return JSONResponse(status_code=400, content={"message": f"Error: No se pudo extraer texto de '{file.filename}'."})
    
    # Elimina los fragmentos antiguos si el archivo ya había sido subido antes
    doc_chunks = [c for c in doc_chunks if c.get("source") != file.filename]
    
    # Procesa el nuevo contenido
    new_chunks_count = 0
    for chunk_text in [text[j:j+1000] for j in range(0, len(text), 1000)]:
        try:
            vec = emb_client.embeddings.create(model=EMBED_MODEL, input=chunk_text).data[0].embedding
            doc_chunks.append({"text": chunk_text, "embedding": np.array(vec), "source": file.filename, "source_id": file.filename})
            new_chunks_count += 1
        except Exception as e: return JSONResponse(status_code=500, content={"message": "Error interno al procesar el archivo."})
        
    # Registra el archivo subido en el log y guarda los cambios
    index_log[file.filename] = {"id": file.filename, "last_modified": time.strftime('%Y-%m-%dT%H:%M:%SZ'), "name": file.filename}
    save_index() # Guarda los cambios en disco
    return {"message": f"'{file.filename}' indexado con éxito con {new_chunks_count} fragmentos."}

@app.get("/reindex", response_class=JSONResponse)
def reindex(folder: str = None, full: bool = False):
    """Activa el proceso de indexación de SharePoint."""
    folder_id = SHAREPOINT_FOLDERS.get(folder.upper()) if folder else None
    if folder and not folder_id: return JSONResponse(status_code=400, content={"message": f"Error: La carpeta '{folder}' no es válida."})
    folder_name = folder.upper() if folder else "Raíz"
    result_message = build_sharepoint_index(folder_id, full_reindex=full)
    mode = "completa" if full else "incremental"
    return {"message": f"Indexación {mode} para '{folder_name}' finalizada. {result_message}"}

# ======================================================
# INICIO DE LA APLICACIÓN
# ======================================================
if __name__ == "__main__":
    # Esta sección solo se ejecuta cuando el script se corre directamente
    # import uvicorn
    # Inicia el servidor web Uvicorn
    # uvicorn.run(app, host="0.0.0.0", port=8000) # Dejado como comentario, ya que load_index() está arriba
    
    # Si prefieres seguir usando `python3 app.py`, descomenta la línea de abajo
    import uvicorn
    uvicorn.run("app:app", host="0.0.0.0", port=8000, reload=True) # Añadido 'reload' para desarrollo